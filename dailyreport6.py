#这个程序不能多次运行，而且中文会出错
#导入所需的包

#还差排水阀排氮阀开启过滤,完成
#稳定时间2s,10s,30s过滤，完成
#控制要求的偏差图，持续超过abs(2)的偏差图
#GUI界面，可选系统号和重跑与否？ 完成
#0.1s采样能改吗？完成
#不同系统的操作条件的更新？未进行


import os
import glob
import gc

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import math


from pptx import Presentation
from pptx.util import Inches,Pt

import tkinter as tk
from tkinter import filedialog,dialog
from tkinter import StringVar
from tkinter import ttk

## 合并和时间稀释
def hebing(folder,signallist,sample_frequency,re_run_flag):

    csv_list = glob.glob('*.csv')#查看同文件夹下的csv文件数
##    if (len(csv_list)>0)&('result.csv' not in csv_list):

    if ((len(csv_list)>0)&(os.path.exists(folder[:(folder.find('CSV'))]+'RESULTS\\')==0) or (re_run_flag==1)):
        for csv in csv_list[:]:
            if ('Systemdata' not in csv) & ('160P' not in csv) & ('639F' not in csv):
                csv_list.remove(csv)
        print(csv_list)
        print(u'共发现%s个CSV文件'% len(csv_list))
        print(u'正在合并............')
        i = 0
        for csv in csv_list: #循环读取同文件夹下的csv文件
            #fr = open(csv,'r').read()
            try:
                fr = pd.read_csv(csv,usecols=signallist,encoding = 'gb2312',skiprows=lambda x: x > 0 and (x-1) % sample_frequency != 0)
            except:
                fr = pd.read_csv(csv,usecols=signallist,skiprows=lambda x: x > 0 and (x-1) % sample_frequency != 0)
##            fr = pd.read_csv(csv,usecols=signallist,encoding = 'gb2312')
            print(len(fr))
            if i==0:
                f = fr
            else:
                f = pd.concat([f,fr],axis =0)
            i = i +1
        f.index = range(len(f))
##        f = f.iloc[range(0,len(f),sample_frequency)]
##        f.index = range(len(f))
    else:
        f = pd.DataFrame(columns = signallist)
    return f


# 选择对应的文件夹
def folder_select(folder,sys_type):
    delete_list = ['txt','xlsx','png','ppt','pptx','csv','docx','rar']
    subfolder_list = glob.glob('*')
    for subfolder in subfolder_list[:]:
        if (sys_type not in subfolder) or (subfolder[(subfolder.find('.')+1):] in delete_list):
            subfolder_list.remove(subfolder)
    ## path supplement

    return subfolder_list

#过滤器

def filters(df,Signal_st,curr_set,temp_set):

    df_index_copy = df.index
    df=df[df[Signal_st]==3]
    if len(df)>0:
        df_id = np.insert(np.diff(df.index),0,[1])
        df_id_break_i = [i for i,x in enumerate(df_id) if x>1]
        df_id_break = df.index[df_id_break_i]
        print(df_id_break)
        if len(df_id_break)==0:
            print('##################only one continous runnnig stage in the data.##############')
            df_id_break_st = [df.index[0]]
            df_id_break_end = [df.index[-1]]
        else:
            df_id_break_st = np.insert(df_id_break,0,df.index[0])
            df_id_break_end_i = [i-1 for i in df_id_break_i]
            df_id_break_end = df.index[df_id_break_end_i]
            df_id_break_end = df_id_break_end.append(df.index[[-1]])
        print(df_id_break_st,'df_id_break_st')
        print(df_id_break_end,'df_id_break_end')
        df_id_warmed = list()
    ##    df_warmed = pd.DataFrame(columns = signallist)
        for k in range(len(df_id_break_st)):
            print(k)
            print(df_id_break_st[k])
            print(df_id_break_end[k])
            i=df_id_break_st[k]


            count = 0
            while (i < df_id_break_end[k]) & (count <100):
    ##            print(i)#10*sample_time
                if (df['CDC_T_TempIn'][i])>(temp_set-2):
                    count = count + 1
                i = i + 1
            print(range(i,df_id_break_end[k]+1))
            df_id_warmed.extend(list(df_index_copy[range(i,df_id_break_end[k]+1)]))
            
        


        df_warmed = df.loc[df_id_warmed]
        df_rated = df_warmed[(df_warmed['SPM_I_StkCurr']<(curr_set+5))&(df_warmed['SPM_I_StkCurr']>(curr_set-5))&(df_warmed['CDC_T_TempIn']>(temp_set-2))&(df_warmed['CDC_T_TempIn']<(temp_set+2))]
        df_id_rated = df_rated.index
        del df_warmed,df_rated
        gc.collect()
    else:
        df_id_break_st = []
        df_id_break_end = []
        df_id_warmed = []
        df_id_rated = []



    
    return df_id_break_st,df_id_break_end,df_id_warmed,df_id_rated

def curr_set_stable(df,Signal_st,Signal_currset,curr_set,temp_set,processdata_hz):

    df_id = np.insert(np.diff(df[Signal_currset]),0,[0])
##%%%%%%%%%%%% transient
    df_id_stable = [i for i,x in enumerate(df_id) if x==0]
##    df_id_stable = [i for i,x in enumerate(df_id) if x==0]
    df_id_break_st,df_id_break_end = stable_select(df_id_stable)
##    for (j,k) in (df_id_break_st[:],df_id_break_end[:]):
##        if k-j<20:
##            df_id_break_st.remove(j)
##            df_id_break_end.remove(k)
##%%%%%%%%%%% stable 2s
    df_id_stable_2s = list()
    for i in range(len(df_id_break_st)-1):
        if (df_id_break_end[i]+1-df_id_break_st[i])>(2/processdata_hz):
            df_id_stable_2s.extend(list(range(df_id_break_st[i]+int(2/processdata_hz),df_id_break_end[i]+1)))

            
##%%%%%%%%%%% stable 10s
    df_id_stable_10s = list()
    for i in range(len(df_id_break_st)-1):
        if (df_id_break_end[i]+1-df_id_break_st[i])>(10/processdata_hz):
            df_id_stable_10s.extend(list(range(df_id_break_st[i]+int(10/processdata_hz),df_id_break_end[i]+1)))

##%%%%%%%%%%% stable 30s            
    df_id_stable_30s = list()
    for i in range(len(df_id_break_st)-1):
        if (df_id_break_end[i]+1-df_id_break_st[i])>(30/processdata_hz):
            df_id_stable_30s.extend(list(range(df_id_break_st[i]+int(30/processdata_hz),df_id_break_end[i]+1)))

##%%%%%%%%%%% stable 60s            
    df_id_stable_60s = list()
    for i in range(len(df_id_break_st)-1):
        if (df_id_break_end[i]+1-df_id_break_st[i])>(60/processdata_hz):
            df_id_stable_60s.extend(list(range(df_id_break_st[i]+int(60/processdata_hz),df_id_break_end[i]+1)))

    print(len(df_id_stable),len(df_id_stable_2s),len(df_id_stable_10s),len(df_id_stable_30s),len(df_id_stable_60s))
    print((set(df_id_stable)&set(df_id_stable_2s))==set(df_id_stable_2s))
    print((set(df_id_stable_2s)&set(df_id_stable_10s))==set(df_id_stable_10s))
    print((set(df_id_stable_10s)&set(df_id_stable_30s))==set(df_id_stable_30s))
    print((set(df_id_stable_30s)&set(df_id_stable_60s))==set(df_id_stable_60s))
    return df_id_stable,df_id_break_st,df_id_break_end,df_id_stable_2s,df_id_stable_10s,df_id_stable_30s,df_id_stable_60s


def stable_select(a_list):
    df_id = np.insert(np.diff(a_list),0,[1])
    df_id_break = [i for i,x in enumerate(df_id) if x>1]
    print(df_id_break)
    if len(df_id_break)==0:
        df_id_break_st = [a_list[0]]
        df_id_break_end = [a_list[-1]]
    else:
        df_id_break_st=[a_list[0]]
        df_id_break_end = list()
        for i in df_id_break:
            df_id_break_st.append(a_list[i])
            df_id_break_end.append(a_list[i-1])
            
##        df_id_break_st = np.insert(df_id_break,0,a_list[0])
##        df_id_break_end = [i-1 for i in df_id_break]
        df_id_break_end.append(a_list[-1])

##    print(df_id_break_st)
##    print(df_id_break_end)
    return df_id_break_st,df_id_break_end


def start_generation():

    if sys_type == "211E-DL3-2-003":
        cell_num = 402
        curr_set = 550
        temp_set = 60
        rawdata_hz = 0.1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_SubSt'
        Signal_st = 'STM_n_Sts'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_PrsOut'
        Signal_testbench = ['Time','AbsTime','Average_1','Average_2']
        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type =="211P-DL3-1-001":
        cell_num = 450
        curr_set = 425
        temp_set = 55
        rawdata_hz = 0.1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_MainSt'
        Signal_st = 'STM_n_StsForVCU'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_HRBPrsIn'
        Signal_testbench = ['Time','AbsTime','Average_1','Average_2']

        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type == "211P-PP-001":
        cell_num = 450
        curr_set = 500
        temp_set = 55
        rawdata_hz = 0.1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_MainSt'
        Signal_st = 'STM_n_StsForVCU'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_HRBPrsIn'
        Signal_testbench = ['Time','AbsTime','Average_1','Average_2']
        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type == "160P-DL3-2-006":### 160P试验车
        cell_num = 274
        curr_set = 425
        temp_set = 55
        rawdata_hz = 1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_MainSt'
        Signal_st = 'STM_n_StsForVCU'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_HRBPrsIn'
        Signal_testbench = []
        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type == "180P-VB-003":
        cell_num = 274
        curr_set = 475
        temp_set = 63.5
        rawdata_hz = 0.1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_MainSt'
        Signal_st = 'STM_n_StsForVCU'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_HRBPrsIn'
        Signal_testbench = []
        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type == "211E-VB-003":
        cell_num = 402
        curr_set = 550
        temp_set = 60
        rawdata_hz = 0.1 #s
        processdata_hz = 0.5 #s
        Signal_subst = 'STM_n_SubSt'
        Signal_st = 'STM_n_Sts'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_PrsOut'
        Signal_testbench = ['Time','AbsTime','Average_1','Average_2']
        opsfilepath = root_path + "\\211P Operation condition.csv"
    elif sys_type == "211P-PP-639":
        cell_num = 450
        curr_set = 250
        temp_set = 55
        rawdata_hz = 1 #s
        processdata_hz = 1 #s
        Signal_subst = 'STM_n_MainSt'
        Signal_st = 'STM_n_StsForVCU'
        Signal_currset = 'SPM_I_CurrSet'
        Signal_anode_prs_out = 'HDC_p_HRBPrsIn'
        Signal_testbench = []
        opsfilepath = root_path + "\\211P Operation condition-T83.csv"

    
    ops_crit = pd.read_csv(opsfilepath)
    os.chdir(root_path)
    ## find the right system file folders
    seleted_folders = folder_select(root_path,sys_type)

    sample_frequency = processdata_hz/rawdata_hz
    signallist = ['SPM_I_StkCurr','SPM_U_StkVolt',
                  'CDC_T_TempIn','CDC_T_TempOut','ADC_p_PrsIn','ADC_dM_MFM','HDC_p_PrsIn',
                  'CDC_N_WCPSpd','HDC_N_HRBSpd','ADC_N_ACPSpd','HDC_b_WDVCmd','HDC_b_NPVCmd','SDM_n_FltGrd']#,'STM_n_StsForVCU'
    signallist.append(Signal_st)
    signallist.append(Signal_anode_prs_out)

    signallist.extend(Signal_testbench)
    signallist.append(Signal_currset)


    
    for folder in seleted_folders:
        os.chdir(root_path + '\\'+ folder + '\\')
        df_rateddata = pd.DataFrame(columns=signallist)
        accum_time = 0
        datefolder_list = glob.glob('*')##find datefile 20210519
        for datefolder in datefolder_list:
            ## try if CSV spelling is right or not, in case tester upload floders with wrong name
            try:
                datefolder_compact = root_path + '\\'+ folder + '\\' + datefolder+'\\CSV\\'
                os.chdir(datefolder_compact)
            except FileNotFoundError:
                continue
            
            df = hebing(datefolder_compact,signallist,sample_frequency,re_run_flag)
            if len(df)>0:
                df['RunTime'] = df['SPM_I_StkCurr']>=5
                print(len(df['SPM_I_StkCurr']))
                print(len(df['RunTime']))
                df['RunTime'] = df['RunTime'].cumsum(0)+accum_time
                accum_time = max(df['RunTime'])
            else:
                continue
            
            df_rateddata = pd.concat([df_rateddata,df[(df['SPM_I_StkCurr']<(curr_set+10))&(df['SPM_I_StkCurr']>(curr_set-10))&(df['CDC_T_TempIn']>(temp_set-2))]],axis =0)
            print(len(df_rateddata),'the length of rated points')


            
            if len(df)>0:

                ## 筛选
                ## 热机完成
                df_id_break_st,df_id_break_end,df_id_warmed,df_id_rated = filters(df,Signal_st,curr_set,temp_set)
                ## 没有purge
                df_id_nopurge = list(df.index[(df['HDC_b_WDVCmd']==0)&(df['HDC_b_NPVCmd']==0)])
                ## 设定稳定后多少s
                df_id_stable,df_id_break_st,df_id_break_end,df_id_stable_2s,df_id_stable_10s,df_id_stable_30s,df_id_stable_60s = curr_set_stable(df,Signal_st,Signal_currset,curr_set,temp_set,processdata_hz)
                ## 没有限功率
                df_id_nofault = list(df.index[(df['SDM_n_FltGrd']==4)|(df['SDM_n_FltGrd']==0)])

                print(set(df_id_stable_60s)==(set(df_id_stable_60s)&set(df_id_stable_30s)))
                print(set(df_id_stable_30s)==(set(df_id_stable_30s)&set(df_id_stable_10s)))
                print(set(df_id_stable_10s)==(set(df_id_stable_2s)&set(df_id_stable_10s)))
                print(set(df_id_stable_2s)==(set(df_id_stable)&set(df_id_stable_2s)))
                


## Stable set ##################
                df_id_filtered = set(df_id_warmed)&set(df_id_stable)

                df_id_filtered_2s = set(df_id_warmed)&set(df_id_stable_2s)
                df_id_filtered_10s = set(df_id_warmed)&set(df_id_stable_10s)
                df_id_filtered_30s = set(df_id_warmed)&set(df_id_stable_30s)
                df_id_filtered_60s = set(df_id_warmed)&set(df_id_stable_60s)
## Stable set ##################
                ## 主要是氢压用
                df_id_filtered_nopurge = set(df_id_warmed)&set(df_id_nopurge)&set(df_id_stable)

                df_id_filtered_nopurge_2s = set(df_id_warmed)&set(df_id_nopurge)&set(df_id_stable_2s)
                df_id_filtered_nopurge_10s = set(df_id_warmed)&set(df_id_nopurge)&set(df_id_stable_10s)
                df_id_filtered_nopurge_30s = set(df_id_warmed)&set(df_id_nopurge)&set(df_id_stable_30s)
                df_id_filtered_nopurge_60s = set(df_id_warmed)&set(df_id_nopurge)&set(df_id_stable_60s)
## Stable set ##################
                ## 主要是水温用
                df_id_filtered_nofault = set(df_id_warmed)&set(df_id_nofault)&set(df_id_stable)
                
##                df_id_filtered_nofault_2s = set(df_id_warmed)&set(df_id_nofault)&set(df_id_stable_2s)
                df_id_filtered_nofault_10s = set(df_id_warmed)&set(df_id_nofault)&set(df_id_stable_10s)
                df_id_filtered_nofault_30s = set(df_id_warmed)&set(df_id_nofault)&set(df_id_stable_30s)
                df_id_filtered_nofault_60s = set(df_id_warmed)&set(df_id_nofault)&set(df_id_stable_60s)

                try:
                    os.mkdir(root_path + '\\'+ folder + '\\' + datefolder+'\\RESULTS\\')
                except:
                    print(datefolder+' RESULTS filefolder already exists.')
                
                os.chdir(root_path + '\\'+ folder + '\\' + datefolder+'\\RESULTS\\')


##                print(len(df_id_stable))
##                print(len(df_id_stable_2s))
##                print(len(df_id_stable_10s))
##                print(len(df_id_stable_30s))
##                print(len(df_id_stable_60s))
                #ret2= list(set(a) | set(b)) #并
                #ret2=list(set(a)-set(b))#差
                #ret2= list(set(a) & set(b)) #交
## Gererate output datafile #############
##                df.loc[df_id_rated].to_csv('./result.csv')

                df.loc[df_id_stable].to_csv('./'+ datefolder + '-stable.csv')

## Gererate output pictures #############                
############################################1
##                fig1, ax1 = plt.subplots()
                fig1 = plt.figure(figsize=(12,5))
                ax1 = fig1.add_subplot(1, 1, 1)
                plt.scatter(df.loc[df_id_rated].index*processdata_hz,df.loc[df_id_rated]['SPM_U_StkVolt']/cell_num*1000,label='SPM_U_StkVolt',s=4)
                ax1.legend()
                ax1.set_title('Average cell voltage vs. Time')
                ax1.set_xlabel('Time')
                ax1.set_ylabel('Average cell voltage')
                ax1.set_xlim([0,500])
                ax1.set_ylim([500,875])
                ax1.grid()
                img_path1 = 'Rated_power_cell_voltage_trend_1.png'
                fig1.savefig(img_path1, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig1, ax1
                gc.collect()

############################################2
                fig2, ax2 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_SPM_U_StkVolt']/cell_num*1000,label='Ops',marker='o')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['SPM_U_StkVolt']/cell_num*1000,label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['SPM_U_StkVolt']/cell_num*1000,label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['SPM_U_StkVolt']/cell_num*1000,label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['SPM_U_StkVolt']/cell_num*1000,label='Stable_30s',s=4)
                
                ax2.legend()
                ax2.set_title('Average Cell voltage vs. Stack current')
                ax2.set_xlabel('Stack current/A')
                ax2.set_ylabel('Cell voltage/mV')
                ax2.set_xlim([0,500])
                ax2.set_ylim([500,900])
                ax2.grid()
                img_path2 = 'Cell_voltage_vs_Stack_current_2.png'
                fig2.savefig(img_path2, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig2, ax2
                gc.collect()

                if ((sys_type != '160P-DL3-2-006')&(sys_type != '180P-VB-003')&(sys_type != '211P-PP-639')):
    ############################################31
                    fig31, ax31 = plt.subplots()
                    plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_SPM_U_StkVolt']/cell_num*1000,label='Ops',marker='o')
                    plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['Average_1'],label='Transient',s=4)
                    plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['Average_1'],label='Stable_2s',s=4)
                    plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['Average_1'],label='Stable_10s',s=4)
                    plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['Average_1'],label='Stable_30s',s=4)
                    
                    ax31.legend()
                    ax31.set_title('Cell voltage vs. Stack current')
                    ax31.set_xlabel('Stack current/A')
                    ax31.set_ylabel('Upper Cell voltage/mV')
                    ax31.set_xlim([0,500])
                    ax31.set_ylim([500,900])
                    ax31.grid()
                    img_path31 = 'Upper_cell_voltage_vs_Stack_current_31.png'
                    fig31.savefig(img_path31, bbox_inches='tight', pad_inches=0, dpi=600)
                    
                    del fig31, ax31
                    gc.collect()
    ############################################32
                    fig32, ax32 = plt.subplots()
                    plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_SPM_U_StkVolt']/cell_num*1000,label='Ops',marker='o')
                    plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['Average_2'],label='Transient',s=4)
                    plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['Average_2'],label='Stable_2s',s=4)
                    plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['Average_2'],label='Stable_10s',s=4)
                    plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['Average_2'],label='Stable_30s',s=4)
                    
                    ax32.legend()
                    ax32.set_title('Cell voltage vs. Stack current')
                    ax32.set_xlabel('Stack current/A')
                    ax32.set_ylabel('Lower Cell voltage/mV')
                    ax32.set_xlim([0,500])
                    ax32.set_ylim([500,900])
                    ax32.grid()
                    img_path32 = 'Lower_cell_voltage_vs_Stack_current_32.png'
                    fig32.savefig(img_path32, bbox_inches='tight', pad_inches=0, dpi=600)
                    
                    del fig32, ax32
                    gc.collect()
############################################41
                fig41, ax41 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_ADC_p_PrsIn'],label='Ops',marker='o')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_ADC_p_PrsIn_UL'],label='Ops_UL',marker='.')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_ADC_p_PrsIn_LL'],label='Ops_LL',marker='.')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['ADC_p_PrsIn'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['ADC_p_PrsIn'],label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['ADC_p_PrsIn'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['ADC_p_PrsIn'],label='Stable_30s',s=4)

                ax41.legend()
                ax41.set_title('Air inlet pressure vs. Stack current')
                ax41.set_xlabel('Stack current/A')
                ax41.set_ylabel('Air inlet pressure/kPa')
                ax41.set_xlim([0,500])
##                ax41.set_ylim([100,230])
                ax41.set_ylim([100,math.ceil(max(ops_crit['ops_ADC_p_PrsIn_UL'])/10)*10])
                ax41.grid()
                img_path41 = 'Air_inlet_pressure_vs_Stack_current_41.png'
                fig41.savefig(img_path41, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig41, ax41
                gc.collect()
############################################42
                fig42, ax42 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_ADC_dM_MFM'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['ADC_dM_MFM'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['ADC_dM_MFM'],label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['ADC_dM_MFM'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['ADC_dM_MFM'],label='Stable_30s',s=4)

                ax42.legend()
                ax42.set_title('Air mass flow vs. Stack current')
                ax42.set_xlabel('Stack current/A')
                ax42.set_ylabel('Air mass flow(g/s)')
                ax42.set_xlim([0,500])
##                ax42.set_ylim([10,140])
                ax42.set_ylim([10,math.ceil(max(ops_crit['ops_ADC_dM_MFM'])*1.2/10)*10])
                ax42.grid()
                img_path42 = 'Air_mass_flow_vs_Stack_current_42.png'
                fig42.savefig(img_path42, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig42, ax42
                gc.collect()
############################################51
                fig51, ax51 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['Anode_dp'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_filtered_nopurge]['SPM_I_StkCurr'],df.loc[df_id_filtered_nopurge]['HDC_p_PrsIn']-df.loc[df_id_filtered_nopurge][Signal_anode_prs_out],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_nopurge_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nopurge_2s]['HDC_p_PrsIn']-df.loc[df_id_filtered_nopurge_2s][Signal_anode_prs_out],label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_nopurge_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nopurge_10s]['HDC_p_PrsIn']-df.loc[df_id_filtered_nopurge_10s][Signal_anode_prs_out],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_nopurge_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nopurge_30s]['HDC_p_PrsIn']-df.loc[df_id_filtered_nopurge_30s][Signal_anode_prs_out],label='Stable_30s',s=4)

                ax51.legend()
                ax51.set_title('Anode dp without purge vs. Stack current')
                ax51.set_xlabel('Stack current/A')
                ax51.set_ylabel('Anode dp/kPa')
                ax51.set_xlim([0,500])
##                ax51.set_ylim([0,60])
                ax51.set_ylim([0,math.ceil(max(ops_crit['Anode_dp'])/10)*10])
                ax51.grid()
                img_path51 = 'Anode_dp_nopurge_51.png'
                fig51.savefig(img_path51, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig51, ax51
                gc.collect()
############################################52
                fig52, ax52 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['Anode_dp'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_filtered]['SPM_I_StkCurr'],df.loc[df_id_filtered]['HDC_p_PrsIn']-df.loc[df_id_filtered][Signal_anode_prs_out],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['HDC_p_PrsIn']-df.loc[df_id_filtered_2s][Signal_anode_prs_out],label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['HDC_p_PrsIn']-df.loc[df_id_filtered_10s][Signal_anode_prs_out],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['HDC_p_PrsIn']-df.loc[df_id_filtered_30s][Signal_anode_prs_out],label='Stable_30s',s=4)

                ax52.legend()
                ax52.set_title('Anode dp with purge vs. Stack current')
                ax52.set_xlabel('Stack current/A')
                ax52.set_ylabel('Anode dp/kPa')
                ax52.set_xlim([0,500])
##                ax52.set_ylim([0,60])
                ax52.set_ylim([0,math.ceil(max(ops_crit['Anode_dp'])/10)*10])
                ax52.grid()
                img_path52 = 'Anode_dp_52.png'
                fig52.savefig(img_path52, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig52, ax52
                gc.collect()
############################################71
                fig71, ax71 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_CDC_T_TempIn'],label='Ops',marker='o')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_CDC_T_TempIn_UL'],label='Ops_UL',marker='.')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_CDC_T_TempIn_LL'],label='Ops_UL',marker='.')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['CDC_T_TempIn'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['CDC_T_TempIn'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['CDC_T_TempIn'],label='Stable_30s',s=4)
                plt.scatter(df.loc[df_id_filtered_60s]['SPM_I_StkCurr'],df.loc[df_id_filtered_60s]['CDC_T_TempIn'],label='Stable_60s',s=4)

                ax71.legend()
                ax71.set_title('Coolant inlet temperature vs. Stack current')
                ax71.set_xlabel('Stack current/A')
                ax71.set_ylabel('Coolant inlet temperature/C')
                ax71.set_xlim([0,500])
##                ax71.set_ylim([20,80])
                ax71.set_ylim([20,math.ceil(max(ops_crit['ops_CDC_T_TempIn_UL'])*1.4/10)*10])
                ax71.grid()
                img_path71 = 'Coolant_in_temp_71.png'
                fig71.savefig(img_path71, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig71, ax71
                gc.collect()
############################################72
                fig72, ax72 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['dT'],label='Ops',marker='o')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['dT_UL'],label='Ops_UL',marker='.')
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['dT_LL'],label='Ops_LL',marker='.')
                plt.scatter(df.loc[df_id_filtered_nofault]['SPM_I_StkCurr'],df.loc[df_id_filtered_nofault]['CDC_T_TempOut']-df.loc[df_id_filtered_nofault]['CDC_T_TempIn'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_nofault_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nofault_10s]['CDC_T_TempOut']-df.loc[df_id_filtered_nofault_10s]['CDC_T_TempIn'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_nofault_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nofault_30s]['CDC_T_TempOut']-df.loc[df_id_filtered_nofault_30s]['CDC_T_TempIn'],label='Stable_30s',s=4)
                plt.scatter(df.loc[df_id_filtered_nofault_60s]['SPM_I_StkCurr'],df.loc[df_id_filtered_nofault_60s]['CDC_T_TempOut']-df.loc[df_id_filtered_nofault_60s]['CDC_T_TempIn'],label='Stable_60s',s=4)

                ax72.legend()
                ax72.set_title('Coolant dT vs. Stack current')
                ax72.set_xlabel('Stack current/A')
                ax72.set_ylabel('Coolant dT/C')
                ax72.set_xlim([0,500])
##                ax72.set_ylim([-5,35])
                ax72.set_ylim([math.floor(min(ops_crit['dT_LL'])/10)*10,math.ceil(max(ops_crit['dT_UL'])/10)*10])
                ax72.grid()
                img_path72 = 'Coolant_dT_72.png'
                fig72.savefig(img_path72, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig72, ax72
                gc.collect()
############################################81
                fig81, ax81 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_HDC_N_HRBSpd'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['HDC_N_HRBSpd'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['HDC_N_HRBSpd'],label='Stable_2s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['HDC_N_HRBSpd'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['HDC_N_HRBSpd'],label='Stable_30s',s=4)

                ax81.legend()
                ax81.set_title('HRB speed vs. Stack current')
                ax81.set_xlabel('Stack current/A')
                ax81.set_ylabel('HRB speed/rpm')
                ax81.set_xlim([0,500])
##                ax1.set_ylim([4000,8000])
                ax81.set_ylim([math.floor(min(ops_crit['ops_HDC_N_HRBSpd'])/1000)*1000,math.ceil(max(ops_crit['ops_HDC_N_HRBSpd'])/1000)*1000])
                ax81.grid()

                img_path81 = 'HRB_speed_81.png'
                fig81.savefig(img_path81, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig81, ax81
                gc.collect()
############################################82
                fig82, ax82 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_CDC_N_WCPSpd'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['CDC_N_WCPSpd'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['CDC_N_WCPSpd'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['CDC_N_WCPSpd'],label='Stable_30s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['CDC_N_WCPSpd'],label='Stable_60s',s=4)

                ax82.legend()
                ax82.set_title('WCP speed vs. Stack current')
                ax82.set_xlabel('Stack current/A')
                ax82.set_ylabel('WCP speed/rpm')
                ax82.set_xlim([0,500])
##                ax82.set_ylim([1000,6000])
                ax82.set_ylim([math.floor(min(ops_crit['ops_CDC_N_WCPSpd'])/1000)*1000,math.ceil(max(ops_crit['ops_CDC_N_WCPSpd'])/1000)*1000])
                ax82.grid()
                img_path82 = 'WCP_speed_82.png'
                fig82.savefig(img_path82, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig82, ax82
                gc.collect()
############################################83
                fig83, ax83 = plt.subplots()
                plt.plot(ops_crit['ops_SPM_I_StkCurr'],ops_crit['ops_ADC_N_ACPSpd'],label='Ops',marker='o')
                plt.scatter(df.loc[df_id_warmed]['SPM_I_StkCurr'],df.loc[df_id_warmed]['ADC_N_ACPSpd'],label='Transient',s=4)
                plt.scatter(df.loc[df_id_filtered_2s]['SPM_I_StkCurr'],df.loc[df_id_filtered_2s]['ADC_N_ACPSpd'],label='Stable_10s',s=4)
                plt.scatter(df.loc[df_id_filtered_10s]['SPM_I_StkCurr'],df.loc[df_id_filtered_10s]['ADC_N_ACPSpd'],label='Stable_30s',s=4)
                plt.scatter(df.loc[df_id_filtered_30s]['SPM_I_StkCurr'],df.loc[df_id_filtered_30s]['ADC_N_ACPSpd'],label='Stable_60s',s=4)

                ax83.legend()
                ax83.set_title('ACP speed vs. Stack current')
                ax83.set_xlabel('Stack current/A')
                ax83.set_ylabel('ACP speed/rpm')
                ax83.set_xlim([0,500])
##                ax83.set_ylim([10000,90000])
                ax83.set_ylim([math.floor(min(ops_crit['ops_ADC_N_ACPSpd'])*0.8/1000)*1000,math.ceil(max(ops_crit['ops_ADC_N_ACPSpd'])*1.2/1000)*1000])
                ax83.grid()
##                ax1.set_xlim([0,500])
                img_path83 = 'ACP_speed_83.png'
                fig83.savefig(img_path83, bbox_inches='tight', pad_inches=0, dpi=600)
                
                del fig83, ax83
                gc.collect()
############################################
## generate pptx on the basis of template
                plt.close('all')
                pptx_path = "211P-DL3.pptx"
                prs = Presentation(root_path+pptx_path)
                
##                prs = Presentation()
##                slide_size = (16,9)
##                prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(slide_size[1])
                
##                prs.slide_height = Inches(9)
##                prs.slide_weight = Inches(16)
                blank_slide_layout = prs.slide_layouts[6]

                slide1 = prs.slides.add_slide(blank_slide_layout)
                slide2 = prs.slides.add_slide(blank_slide_layout)
                if ((sys_type != '160P-DL3-2-006')&(sys_type != '180P-VB-003')):
                    slide3 = prs.slides.add_slide(blank_slide_layout)

                slide4 = prs.slides.add_slide(blank_slide_layout)

                slide5 = prs.slides.add_slide(blank_slide_layout)

                slide7 = prs.slides.add_slide(blank_slide_layout)

                slide8 = prs.slides.add_slide(blank_slide_layout)

                
##                left = top = Inches(1)
## On-screen Show (16:9) has slide dimensions of 10 inches x 5.625 inches.

                new_paragraph1 = slide1.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph1.paragraphs[0].text = '系统整体电压变化'
                new_paragraph1.paragraphs[0].font.size = Pt(20)
                slide1.shapes.add_picture(img_path1, left=Inches(4), top=Inches(2))
                
                new_paragraph2 = slide2.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph2.paragraphs[0].text = '系统整体电压分布'
                new_paragraph2.paragraphs[0].font.size = Pt(20)
                slide2.shapes.add_picture(img_path2, left=Inches(4), top=Inches(2))

                if ((sys_type != '160P-DL3-2-006')&(sys_type != '180P-VB-003')&(sys_type != '211P-PP-639')):
                    new_paragraph3 = slide3.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                    new_paragraph3.paragraphs[0].text = '上下堆整体电压分布'
                    new_paragraph3.paragraphs[0].font.size = Pt(20)
                    slide3.shapes.add_picture(img_path31, left=Inches(0.5), top=Inches(2))
                    slide3.shapes.add_picture(img_path32, left=Inches(7), top=Inches(2))

                new_paragraph4 = slide4.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph4.paragraphs[0].text = '操作条件——阴极过量系数和压力'
                new_paragraph4.paragraphs[0].font.size = Pt(20)
                slide4.shapes.add_picture(img_path41, left=Inches(0.5), top=Inches(2))
                slide4.shapes.add_picture(img_path42, left=Inches(7.5), top=Inches(2))

                new_paragraph5 = slide5.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph5.paragraphs[0].text = '操作条件——阳极压差'
                new_paragraph5.paragraphs[0].font.size = Pt(20)
                slide5.shapes.add_picture(img_path51, left=Inches(0.5), top=Inches(2))
                slide5.shapes.add_picture(img_path52, left=Inches(7.5), top=Inches(2))

                new_paragraph7 = slide7.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph7.paragraphs[0].text = '操作条件——入口水温和温差'
                new_paragraph7.paragraphs[0].font.size = Pt(20)
                slide7.shapes.add_picture(img_path71, left=Inches(0.5), top=Inches(2))
                slide7.shapes.add_picture(img_path72, left=Inches(7.5), top=Inches(2))

                new_paragraph8 = slide8.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1)).text_frame
                new_paragraph8.paragraphs[0].text = '操作条件——泵使用转速分布'
                new_paragraph8.paragraphs[0].font.size = Pt(20)
                slide8.shapes.add_picture(img_path81, left=Inches(7.5), top=Inches(3))
                slide8.shapes.add_picture(img_path82, left=Inches(0), top=Inches(3))
                slide8.shapes.add_picture(img_path83, left=Inches(6), top=Inches(0))
                prs.save(folder+'-'+datefolder+'.pptx')
                ## slide3 and new_paragraph3 is optional, so do not delete it in case of errors!!
                del prs,new_paragraph1,new_paragraph2,new_paragraph4,new_paragraph5,new_paragraph7,new_paragraph8,slide1,slide2,slide4,slide5,slide7,slide8

            else:
                continue

            
        os.chdir(root_path + '\\'+ folder + '\\')
        df_rateddata.to_csv('./'+ folder + '-rateddata.csv')

    return 0


def ini_generation():
    global re_run_flag,sys_type
    re_run_flag = CheckVar1.get()
    sys_type = comboxlist.get()
    print('Now you want to process:',sys_type)
    print('And you also want to re_run historical data(Yes = 1, No = 0):',re_run_flag)
    start_generation()
    print('All analysis completed, please choose again or Quit GUI!')
    return 0

def open_file():
    '''
    打开文件
    :return:
    '''

    global root_path

    root_path = filedialog.askdirectory(title=u'选择路径', initialdir='D:\\2.DUR&EMS\\2.DURDAILYREPORT')
    root_path = root_path+'\\'
    print('选择的数据路径为：', root_path)



##主程序

if __name__ == '__main__':

##    root_path = "Z:\\2.耐久试验\\"
##    root_path = "D:\\2.DUR&EMS\\2.DURDAILYREPORT\\Trial\\"
##    root_path = "D:\\2.DUR&EMS\\2.DURDAILYREPORT\\"
##    sys_type = "211P-DL3-1-001"
##    sample_frequency = 5
    root_path = 'D:\\2.DUR&EMS\\2.DURDAILYREPORT\\'

## windows loop
    window = tk.Tk()
    window.title('DataProcessGUI')  # 标题
    window.geometry('200x400')  # 窗口尺寸

    comvalue=tk.StringVar()#窗体自带的文本，新建一个值
    comboxlist=ttk.Combobox(window,textvariable=comvalue) #初始化
    comboxlist["values"]=("211P-DL3-1-001","211E-DL3-2-003","211P-PP-001","160P-DL3-2-006","180P-VB-003","211E-VB-003","211P-PP-639")
    comboxlist.current(0) #选择第一个
    comboxlist.grid(row=1,column=0)
    sys_type = comboxlist.get()

##    sys_type = "211P-DL3-1-001"


    CheckVar1 = tk.IntVar()
    checkbutton = tk.Checkbutton(window, text = "是否重新运行历史数据?", variable = CheckVar1, \
                 onvalue = 1, offvalue = 0, height=5, \
                 width = 20)
    checkbutton.grid(row=2,column=0)
    re_run_flag = CheckVar1.get()

    button = tk.Button(window, text='运行', width=10, height=4, command=ini_generation)
    button.grid(row=3,column=0)

    button2 = tk.Button(window, text='选择路径', width=10, height=4, command=open_file)
    button2.grid(row=4,column=0)
        
    
##    signallist = ['STM_n_MainSt','SPM_I_StkCurr','SPM_U_StkVolt','UDS_Id_CurrDensitySet']
##    f = hebing(root_path,signallist)# 3.contentate all the data to on file
    window.mainloop()  # 显示
    os.chdir(root_path)

